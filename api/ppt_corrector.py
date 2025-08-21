#!/usr/bin/env python3
"""
PowerPoint 자동 수정 모듈 (API용)
"""

import os
import shutil
from pptx import Presentation
from typing import List, Dict
from datetime import datetime

class PPTAutoCorrector:
    def __init__(self, pptx_file: str):
        self.pptx_file = pptx_file
        self.presentation = Presentation(pptx_file)
        self.corrections = []
        self.backup_file = None
        
    def create_backup(self) -> str:
        """원본 파일 백업"""
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        backup_name = f"{os.path.splitext(self.pptx_file)[0]}_backup_{timestamp}.pptx"
        shutil.copy2(self.pptx_file, backup_name)
        self.backup_file = backup_name
        return backup_name
    
    def find_and_replace_in_shape(self, shape, old_text: str, new_text: str) -> bool:
        """도형 내 텍스트 찾아서 교체 (서식 유지)"""
        replaced = False
        
        if not hasattr(shape, 'text_frame'):
            return False
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    # 원본 서식 저장
                    original_font_name = run.font.name
                    original_font_size = run.font.size
                    original_bold = run.font.bold
                    original_italic = run.font.italic
                    original_color = run.font.color
                    
                    # 텍스트 교체
                    run.text = run.text.replace(old_text, new_text)
                    
                    # 서식 복원 (None이 아닌 경우만)
                    if original_font_name:
                        run.font.name = original_font_name
                    if original_font_size:
                        run.font.size = original_font_size
                    if original_bold is not None:
                        run.font.bold = original_bold
                    if original_italic is not None:
                        run.font.italic = original_italic
                        
                    replaced = True
                    
        return replaced
    
    def apply_corrections(self, corrections_list: List[Dict]) -> Dict:
        """
        선택된 맞춤법 오류들을 PPT 파일에 직접 적용
        
        corrections_list: [
            {
                'slide_number': 1,
                'original': '됬습니다',
                'corrected': '됐습니다',
                'apply': True  # 사용자가 선택한 것만
            }
        ]
        """
        applied_count = 0
        failed_count = 0
        results = []
        
        for correction in corrections_list:
            if not correction.get('apply', False):
                continue
                
            slide_num = correction['slide_number']
            original = correction['original']
            corrected = correction['corrected']
            
            try:
                slide = self.presentation.slides[slide_num - 1]
                slide_replaced = False
                
                # 모든 도형에서 텍스트 교체
                for shape in slide.shapes:
                    if hasattr(shape, "text") and original in shape.text:
                        if self.find_and_replace_in_shape(shape, original, corrected):
                            slide_replaced = True
                    
                    # 테이블 내 텍스트도 교체
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if original in cell.text:
                                    # 테이블 셀의 텍스트 프레임 처리
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            if original in run.text:
                                                run.text = run.text.replace(original, corrected)
                                                slide_replaced = True
                
                if slide_replaced:
                    applied_count += 1
                    results.append({
                        'slide': slide_num,
                        'original': original,
                        'corrected': corrected,
                        'status': 'success'
                    })
                else:
                    failed_count += 1
                    results.append({
                        'slide': slide_num,
                        'original': original,
                        'corrected': corrected,
                        'status': 'not_found'
                    })
                    
            except Exception as e:
                failed_count += 1
                results.append({
                    'slide': slide_num,
                    'original': original,
                    'corrected': corrected,
                    'status': 'error',
                    'error': str(e)
                })
        
        return {
            'applied': applied_count,
            'failed': failed_count,
            'total': len([c for c in corrections_list if c.get('apply', False)]),
            'details': results
        }
    
    def save_corrected_file(self, output_path: str = None) -> str:
        """수정된 파일 저장"""
        if output_path is None:
            base_name = os.path.splitext(self.pptx_file)[0]
            output_path = f"{base_name}_수정됨_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        self.presentation.save(output_path)
        return output_path