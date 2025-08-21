#!/usr/bin/env python3
import os
import json
from openai import OpenAI
from pptx import Presentation
from typing import List, Dict
import argparse
from datetime import datetime
import time

class PPTSpellCheckerGPT:
    def __init__(self, pptx_file: str, api_key: str = None):
        self.pptx_file = pptx_file
        self.presentation = Presentation(pptx_file)
        self.errors = []
        
        # API 키 설정 (환경변수 또는 직접 전달)
        if api_key:
            self.client = OpenAI(api_key=api_key)
        else:
            api_key = os.getenv('OPENAI_API_KEY')
            if not api_key:
                raise ValueError("OpenAI API 키가 필요합니다. OPENAI_API_KEY 환경변수를 설정하거나 --api-key 옵션을 사용하세요.")
            self.client = OpenAI(api_key=api_key)
        
    def extract_text_from_slides(self) -> List[Dict]:
        """슬라이드에서 모든 텍스트 추출"""
        texts = []
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append({
                        'slide_number': slide_num,
                        'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
                        'text': shape.text.strip(),
                        'position': f"Slide {slide_num}"
                    })
                    
                # 테이블 내 텍스트도 추출
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                texts.append({
                                    'slide_number': slide_num,
                                    'shape_name': f"Table[{row_idx}][{col_idx}]",
                                    'text': cell.text.strip(),
                                    'position': f"Slide {slide_num}, Table Cell [{row_idx+1},{col_idx+1}]"
                                })
        
        return texts
    
    def check_spelling_gpt(self, text: str, model: str = "gpt-4o-mini") -> Dict:
        """ChatGPT API를 사용한 맞춤법 검사"""
        
        prompt = f"""다음 한국어 텍스트의 맞춤법과 문법을 검사해주세요. 
오류가 있다면 수정된 텍스트와 함께 JSON 형식으로 응답해주세요.
오류가 없다면 has_errors를 false로 설정해주세요.

텍스트: "{text}"

응답 형식:
{{
    "has_errors": true/false,
    "corrected": "수정된 텍스트",
    "errors": ["발견된 오류1", "발견된 오류2", ...],
    "error_count": 오류 개수
}}

주의사항:
- 반드시 JSON 형식으로만 응답
- 수정된 텍스트는 원본의 의미를 유지하면서 자연스럽게 수정
- 맞춤법, 띄어쓰기, 문법 오류를 모두 검사"""

        try:
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": "당신은 한국어 맞춤법 전문가입니다. JSON 형식으로만 응답합니다."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                response_format={"type": "json_object"}
            )
            
            result = json.loads(response.choices[0].message.content)
            
            return {
                'has_errors': result.get('has_errors', False),
                'original': text,
                'corrected': result.get('corrected', text),
                'error_count': result.get('error_count', 0),
                'errors': result.get('errors', [])
            }
            
        except Exception as e:
            print(f"GPT API 오류: {e}")
            return {
                'has_errors': False,
                'original': text,
                'corrected': text,
                'error_count': 0,
                'error': str(e)
            }
    
    def check_all_slides(self, model: str = "gpt-4o-mini", batch_size: int = 5) -> List[Dict]:
        """모든 슬라이드의 맞춤법 검사 수행"""
        texts = self.extract_text_from_slides()
        results = []
        
        print(f"총 {len(texts)}개의 텍스트 블록을 검사합니다...")
        print(f"사용 모델: {model}")
        
        for idx, text_info in enumerate(texts, 1):
            print(f"검사 중... [{idx}/{len(texts)}] - {text_info['position']}", end='\r')
            
            # 맞춤법 검사
            spell_result = self.check_spelling_gpt(text_info['text'], model)
            
            if spell_result['has_errors']:
                results.append({
                    **text_info,
                    **spell_result
                })
            
            # API 호출 제한을 위한 짧은 대기
            if idx % batch_size == 0:
                time.sleep(1)
        
        print(f"\n검사 완료! 총 {len(results)}개의 오류를 발견했습니다.")
        self.errors = results
        return results
    
    def generate_report(self, output_file: str = None) -> str:
        """검사 결과 리포트 생성"""
        if output_file is None:
            base_name = os.path.splitext(self.pptx_file)[0]
            output_file = f"{base_name}_GPT맞춤법검사_결과_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"=" * 80 + "\n")
            f.write(f"PowerPoint 맞춤법 검사 결과 (ChatGPT)\n")
            f.write(f"파일: {self.pptx_file}\n")
            f.write(f"검사 일시: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"발견된 오류: {len(self.errors)}개\n")
            f.write(f"=" * 80 + "\n\n")
            
            if not self.errors:
                f.write("맞춤법 오류를 발견하지 못했습니다!\n")
            else:
                current_slide = 0
                for error in self.errors:
                    if error['slide_number'] != current_slide:
                        current_slide = error['slide_number']
                        f.write(f"\n[슬라이드 {current_slide}]\n")
                        f.write("-" * 40 + "\n")
                    
                    f.write(f"위치: {error['position']}\n")
                    f.write(f"원본: {error['original']}\n")
                    f.write(f"수정: {error['corrected']}\n")
                    if 'errors' in error and error['errors']:
                        f.write(f"발견된 오류: {', '.join(error['errors'])}\n")
                    f.write(f"오류 개수: {error['error_count']}\n")
                    f.write("\n")
        
        print(f"리포트가 저장되었습니다: {output_file}")
        return output_file
    
    def generate_json_report(self, output_file: str = None) -> str:
        """JSON 형식으로 리포트 생성"""
        if output_file is None:
            base_name = os.path.splitext(self.pptx_file)[0]
            output_file = f"{base_name}_GPT맞춤법검사_결과_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        
        report_data = {
            'file': self.pptx_file,
            'check_date': datetime.now().isoformat(),
            'total_errors': len(self.errors),
            'errors': self.errors
        }
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(report_data, f, ensure_ascii=False, indent=2)
        
        print(f"JSON 리포트가 저장되었습니다: {output_file}")
        return output_file

def main():
    parser = argparse.ArgumentParser(description='ChatGPT를 사용한 PowerPoint 파일 맞춤법 검사')
    parser.add_argument('file', help='검사할 PowerPoint 파일 경로')
    parser.add_argument('--api-key', help='OpenAI API 키 (또는 OPENAI_API_KEY 환경변수 사용)')
    parser.add_argument('--model', default='gpt-4o-mini', 
                       choices=['gpt-4o-mini', 'gpt-4o', 'gpt-3.5-turbo'],
                       help='사용할 GPT 모델 (기본: gpt-4o-mini)')
    parser.add_argument('--json', action='store_true', help='JSON 형식으로 결과 출력')
    parser.add_argument('--output', '-o', help='결과 파일명 지정')
    
    args = parser.parse_args()
    
    if not os.path.exists(args.file):
        print(f"파일을 찾을 수 없습니다: {args.file}")
        return
    
    try:
        # 맞춤법 검사 수행
        checker = PPTSpellCheckerGPT(args.file, api_key=args.api_key)
        checker.check_all_slides(model=args.model)
        
        # 리포트 생성
        if args.json:
            checker.generate_json_report(args.output)
        else:
            checker.generate_report(args.output)
            
    except ValueError as e:
        print(f"오류: {e}")
        print("\n사용법:")
        print("1. 환경변수 설정: export OPENAI_API_KEY='your-api-key'")
        print("2. 또는 직접 전달: python ppt_spell_checker_gpt.py file.pptx --api-key 'your-api-key'")

if __name__ == "__main__":
    main()