#!/usr/bin/env python3
"""
PowerPoint 맞춤법 검사 모듈 (API용)
"""

import os
import json
import requests
from openai import OpenAI
from pptx import Presentation
from typing import List, Dict
import time

class PPTSpellChecker:
    def __init__(self, pptx_file: str, api_key: str = None):
        self.pptx_file = pptx_file
        self.presentation = Presentation(pptx_file)
        self.errors = []
        
        # API 키 설정
        if api_key:
            self.client = OpenAI(api_key=api_key)
        else:
            api_key = os.getenv('OPENAI_API_KEY')
            if not api_key:
                raise ValueError("OpenAI API 키가 필요합니다. OPENAI_API_KEY 환경변수를 설정하세요.")
            self.client = OpenAI(api_key=api_key)
        
    def extract_text_from_slides(self) -> List[Dict]:
        """슬라이드에서 모든 텍스트 추출"""
        texts = []
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append({
                        'slide_number': slide_num,
                        'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
                        'text': shape.text.strip(),
                        'position': f"Slide {slide_num}"
                    })
                    
                # 테이블 내 텍스트도 추출
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                texts.append({
                                    'slide_number': slide_num,
                                    'shape_name': f"Table[{row_idx}][{col_idx}]",
                                    'text': cell.text.strip(),
                                    'position': f"Slide {slide_num}, Table Cell [{row_idx+1},{col_idx+1}]"
                                })
        
        return texts
    
    def check_spelling_simple(self, text: str) -> Dict:
        """간단한 규칙 기반 맞춤법 검사 (백업용)"""
        common_errors = {
            '됬': '됐',
            '됀': '된',
            '어떻개': '어떻게',
            '안됀': '안 된',
            '안되': '안 돼',
            '그래도되': '그래도 돼',
            '하던지': '하든지',
            '어떤지': '어떻든지',
            '웬지': '왠지',
            '왠일': '웬일',
            '갈께': '갈게',
            '할께': '할게',
            '있을께': '있을게',
            '되요': '돼요',
            '되어': '돼',
            '되서': '돼서',
            '않하': '안 하',
            '그런대': '그런데',
            '하세여': '하세요',
            '하십시요': '하십시오',
            '해주세여': '해주세요',
            '몇일': '며칠',
            '몇 일': '며칠',
            '예기': '얘기',
            '그렇치': '그렇지',
            '그치': '그렇지',
            '어떻던': '어떻든',
            '틀리게': '다르게',
            '틀려': '달라',
        }
        
        corrected = text
        errors_found = 0
        
        for error, correction in common_errors.items():
            if error in text:
                corrected = corrected.replace(error, correction)
                errors_found += 1
        
        return {
            'has_errors': errors_found > 0,
            'original': text,
            'corrected': corrected,
            'error_count': errors_found
        }
    
    def check_spelling_gpt(self, text: str, model: str = "gpt-4o-mini") -> Dict:
        """ChatGPT API를 사용한 맞춤법 검사"""
        
        prompt = f"""다음 한국어 텍스트의 맞춤법과 문법을 검사해주세요. 
오류가 있다면 수정된 텍스트와 함께 JSON 형식으로 응답해주세요.
오류가 없다면 has_errors를 false로 설정해주세요.

텍스트: "{text}"

응답 형식:
{{
    "has_errors": true/false,
    "corrected": "수정된 텍스트",
    "errors": ["발견된 오류1", "발견된 오류2", ...],
    "error_count": 오류 개수
}}

주의사항:
- 반드시 JSON 형식으로만 응답
- 수정된 텍스트는 원본의 의미를 유지하면서 자연스럽게 수정
- 맞춤법, 띄어쓰기, 문법 오류를 모두 검사"""

        try:
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": "당신은 한국어 맞춤법 전문가입니다. JSON 형식으로만 응답합니다."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                response_format={"type": "json_object"}
            )
            
            result = json.loads(response.choices[0].message.content)
            
            return {
                'has_errors': result.get('has_errors', False),
                'original': text,
                'corrected': result.get('corrected', text),
                'error_count': result.get('error_count', 0),
                'errors': result.get('errors', [])
            }
            
        except Exception as e:
            print(f"GPT API 오류: {e}")
            # GPT 실패 시 간단한 검사 사용
            return self.check_spelling_simple(text)
    
    def check_all_slides(self, model: str = "gpt-4o-mini", use_gpt: bool = True) -> List[Dict]:
        """모든 슬라이드의 맞춤법 검사 수행"""
        texts = self.extract_text_from_slides()
        results = []
        
        print(f"총 {len(texts)}개의 텍스트 블록을 검사합니다...")
        
        for idx, text_info in enumerate(texts, 1):
            print(f"검사 중... [{idx}/{len(texts)}]", end='\r')
            
            # 맞춤법 검사
            if use_gpt:
                spell_result = self.check_spelling_gpt(text_info['text'], model)
            else:
                spell_result = self.check_spelling_simple(text_info['text'])
            
            if spell_result['has_errors']:
                results.append({
                    **text_info,
                    **spell_result
                })
            
            # API 호출 제한을 위한 짧은 대기
            if use_gpt and idx % 5 == 0:
                time.sleep(0.5)
        
        print(f"\n검사 완료! 총 {len(results)}개의 오류를 발견했습니다.")
        self.errors = results
        return results